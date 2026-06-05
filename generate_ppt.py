from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color palette ──────────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x23, 0x3A)   # dark navy  (background)
C_ACCENT = RGBColor(0x00, 0xB0, 0xF0)   # bright cyan
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xCC, 0xE8, 0xF4)   # light blue text
C_YELLOW = RGBColor(0xFF, 0xD7, 0x00)
C_GREEN  = RGBColor(0x70, 0xE0, 0x7E)
C_GRAY   = RGBColor(0xA0, 0xB0, 0xC0)
C_BOX    = RGBColor(0x22, 0x35, 0x55)   # card background

# ── Helper functions ───────────────────────────────────────────────────────────

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color=C_DARK):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_line(slide, l, t, w, color=C_ACCENT, thickness=Pt(2)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.04))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text(slide, text, l, t, w, h, size=Pt(18), bold=False,
             color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微軟正黑體"
    return txBox

def add_bullet_box(slide, items, l, t, w, h, title=None,
                   title_color=C_ACCENT, item_color=C_WHITE,
                   bg_color=C_BOX, icon="▸"):
    add_rect(slide, l, t, w, h, bg_color)
    y = t + 0.18
    if title:
        add_text(slide, title, l+0.2, y, w-0.3, 0.38,
                 size=Pt(15), bold=True, color=title_color)
        y += 0.40
    for item in items:
        add_text(slide, f"{icon}  {item}", l+0.2, y, w-0.35, 0.34,
                 size=Pt(13), color=item_color)
        y += 0.35

# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT)
add_rect(sl, 0, 7.38, 13.33, 0.12, C_ACCENT)

# Decorative side bar
add_rect(sl, 0, 0, 0.35, 7.5, C_BOX)
add_rect(sl, 0, 2.8, 0.35, 1.9, C_ACCENT)

add_text(sl, "智慧實驗室管理系統", 0.7, 1.6, 12, 1.3,
         size=Pt(46), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "人員管制 × 水耕養殖監控", 0.7, 2.9, 12, 0.7,
         size=Pt(28), bold=False, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_line(sl, 2.5, 3.75, 8.33)

add_text(sl, "YOLOv8  ·  InsightFace ArcFace  ·  Python Flask  ·  IoT Sensors",
         0.7, 4.0, 12, 0.5,
         size=Pt(15), color=C_LIGHT, align=PP_ALIGN.CENTER)
add_text(sl, "2026", 0.7, 5.5, 12, 0.5,
         size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 — 目錄
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT)

add_text(sl, "報告大綱", 0.5, 0.25, 12, 0.65,
         size=Pt(30), bold=True, color=C_WHITE)
add_line(sl, 0.5, 0.95, 12.33)

sections = [
    ("01", "研究背景與動機"),
    ("02", "系統架構與技術"),
    ("03", "數據收集方法"),
    ("04", "製作過程"),
    ("05", "研究成果"),
    ("06", "結論與未來展望"),
]
cols = [(0.5, 1.15), (4.5, 1.15), (8.5, 1.15),
        (0.5, 3.9),  (4.5, 3.9),  (8.5, 3.9)]

for i, ((lx, ty), (num, title)) in enumerate(zip(cols, sections)):
    add_rect(sl, lx, ty, 3.8, 2.4, C_BOX)
    add_rect(sl, lx, ty, 3.8, 0.45, C_ACCENT)
    add_text(sl, num, lx+0.15, ty+0.05, 1, 0.4,
             size=Pt(20), bold=True, color=C_DARK)
    add_text(sl, title, lx+0.15, ty+0.6, 3.5, 1.6,
             size=Pt(18), bold=True, color=C_WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 3 — 研究背景與動機
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT)
add_text(sl, "01  研究背景與動機", 0.5, 0.2, 12, 0.6,
         size=Pt(28), bold=True, color=C_WHITE)
add_line(sl, 0.5, 0.85, 12.33)

problems = [
    "傳統門禁依賴識別卡，易發生代刷或遺失",
    "人工簽到本無法即時查詢在室人員狀況",
    "養殖池水質異常須人員定時巡視，費時費力",
    "人員安全與環境監控各自獨立，管理分散",
]
goals = [
    "無接觸自動化人員管制",
    "整合兩大系統於單一儀表板",
    "即時異常警報，縮短反應時間",
    "驗證低成本 AI 攝影機可行性",
]

add_bullet_box(sl, problems, 0.5, 1.05, 6.1, 5.8,
               title="現有問題", icon="✗", item_color=RGBColor(0xFF,0xA0,0xA0))
add_bullet_box(sl, goals,   6.8, 1.05, 6.1, 5.8,
               title="解決目標", icon="✓", item_color=C_GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 4 — 系統架構
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT)
add_text(sl, "02  系統架構與技術", 0.5, 0.2, 12, 0.6,
         size=Pt(28), bold=True, color=C_WHITE)
add_line(sl, 0.5, 0.85, 12.33)

# Tech stack boxes
tech = [
    ("影像輸入", ["4 支 USB 監控攝影機", "各自獨立 Thread", "5 fps 運行"]),
    ("AI 偵測", ["YOLOv8n 人體偵測", "InsightFace ArcFace", "512D Embedding"]),
    ("後端核心", ["Python Flask", "Waitress 32 threads", "MySQL 資料庫"]),
    ("通知整合", ["Telegram Bot API", "即時推送截圖", "分級警報機制"]),
]
for i, (title, items) in enumerate(tech):
    lx = 0.5 + i * 3.2
    add_rect(sl, lx, 1.1, 3.0, 3.5, C_BOX)
    add_rect(sl, lx, 1.1, 3.0, 0.5, C_ACCENT)
    add_text(sl, title, lx+0.1, 1.12, 2.8, 0.45,
             size=Pt(15), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(sl, f"• {item}", lx+0.15, 1.75+j*0.55, 2.7, 0.5,
                 size=Pt(13), color=C_WHITE)

# Arrow flow
add_text(sl, "攝影機 → AI 偵測 → State Machine → MySQL → 儀表板 → Telegram",
         0.5, 4.85, 12.33, 0.5,
         size=Pt(14), color=C_ACCENT, align=PP_ALIGN.CENTER)
add_line(sl, 0.5, 4.8, 12.33, color=C_BOX)

# Camera roles
cam_data = [
    ("攝影機 1", "概覽 + 人體追蹤"),
    ("攝影機 2", "入口 ENTER 觸發\n陌生人偵測"),
    ("攝影機 3", "室內 Presence 確認"),
    ("攝影機 4", "出口 EXIT 觸發"),
]
for i, (cam, role) in enumerate(cam_data):
    lx = 0.5 + i * 3.2
    add_rect(sl, lx, 5.45, 3.0, 1.75, C_BOX)
    add_text(sl, cam,  lx+0.1, 5.5,  2.8, 0.4, size=Pt(13), bold=True, color=C_YELLOW)
    add_text(sl, role, lx+0.1, 5.9,  2.8, 1.1, size=Pt(12), color=C_LIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 — 數據收集
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT)
add_text(sl, "03  數據收集方法", 0.5, 0.2, 12, 0.6,
         size=Pt(28), bold=True, color=C_WHITE)
add_line(sl, 0.5, 0.85, 12.33)

# Left: enrollment process
steps = [
    "① 開啟人員管理頁面",
    "② 新增成員並填寫資料",
    "③ 選擇攝影機 2 或 3 拍攝",
    "④ 系統顯示即時 Pose Guide",
    "⑤ 依序完成 5 種角度拍攝",
    "⑥ 系統自動重建 Embedding",
]
add_bullet_box(sl, steps, 0.5, 1.05, 5.5, 5.8,
               title="人臉登錄流程 (Face Enrollment)", icon="")

# Right: 5 poses + principle
poses = ["正面", "左側", "右側", "低頭", "抬頭"]
add_text(sl, "5 種拍攝角度", 6.3, 1.05, 6.7, 0.5,
         size=Pt(16), bold=True, color=C_ACCENT)
for i, p in enumerate(poses):
    lx = 6.3 + (i % 3) * 2.2
    ty = 1.65 + (i // 3) * 1.5
    add_rect(sl, lx, ty, 2.0, 1.3, C_BOX)
    add_text(sl, p, lx, ty+0.35, 2.0, 0.55,
             size=Pt(18), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_rect(sl, 6.3, 4.55, 6.7, 2.2, RGBColor(0x15,0x2A,0x1A))
add_text(sl, "關鍵發現", 6.5, 4.6, 6.3, 0.4,
         size=Pt(14), bold=True, color=C_GREEN)
add_text(sl,
    "以手機拍攝訓練照 → 信心值 0.42–0.51（低於門檻）\n"
    "改用監控攝影機現場拍攝 → 信心值提升至 0.63–0.79\n\n"
    "結論：訓練資料須與部署環境一致（Domain Consistency）",
    6.5, 5.05, 6.3, 1.55,
    size=Pt(13), color=C_LIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 6 — 製作過程
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT)
add_text(sl, "04  製作過程", 0.5, 0.2, 12, 0.6,
         size=Pt(28), bold=True, color=C_WHITE)
add_line(sl, 0.5, 0.85, 12.33)

timeline = [
    ("Step 1", "感測器連線\n與資料收集", "連接 IoT Node\n建立 MySQL 資料庫"),
    ("Step 2", "系統架構\n設計", "規劃 API、DB 結構\n前後端分工"),
    ("Step 3", "感測器監控\n網頁", "即時儀表板\nAlert 機制"),
    ("Step 4", "AI 人員管制\n系統", "YOLO + ArcFace\nState Machine"),
    ("Step 5", "測試與\n持續優化", "修正 False EXIT\nUNKNOWN 偵測"),
]

for i, (step, title, desc) in enumerate(timeline):
    lx = 0.4 + i * 2.55
    add_rect(sl, lx, 1.1, 2.35, 0.55, C_ACCENT)
    add_text(sl, step, lx, 1.12, 2.35, 0.5,
             size=Pt(14), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_rect(sl, lx, 1.65, 2.35, 4.8, C_BOX)
    add_text(sl, title, lx+0.1, 1.75, 2.15, 0.9,
             size=Pt(15), bold=True, color=C_WHITE)
    add_text(sl, desc, lx+0.1, 2.75, 2.15, 3.5,
             size=Pt(12), color=C_LIGHT)

# bottom bar: key fixes
add_rect(sl, 0.4, 6.55, 12.53, 0.75, RGBColor(0x12,0x1E,0x35))
add_text(sl,
    "主要修正：Camera backend DSHOW→MSMF  ·  重拍訓練照  ·  "
    "False EXIT grace period  ·  UNKNOWN 獨立偵測邏輯",
    0.5, 6.6, 12.33, 0.6,
    size=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 — 成果（人員管制）
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT)
add_text(sl, "05  研究成果 — 人員管制", 0.5, 0.2, 12, 0.6,
         size=Pt(28), bold=True, color=C_WHITE)
add_line(sl, 0.5, 0.85, 12.33)

# Confidence table
headers = ["成員", "改善前（手機拍攝）", "改善後（攝影機拍攝）"]
rows = [
    ["Tongen", "0.63", "0.63 – 0.77"],
    ["jiabao",  "0.42 – 0.51  ✗", "0.68  ✓"],
    ["TT",      "無法偵測  ✗",     "0.79  ✓"],
]
add_text(sl, "人臉辨識信心值比較", 0.5, 1.0, 7.5, 0.45,
         size=Pt(15), bold=True, color=C_ACCENT)

col_w = [1.8, 2.6, 2.6]
col_x = [0.5, 2.35, 5.0]
row_h = 0.55
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(sl, cx, 1.5, cw-0.05, row_h, C_ACCENT)
    add_text(sl, hdr, cx+0.05, 1.52, cw-0.1, row_h-0.05,
             size=Pt(12), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

row_colors = [C_BOX, RGBColor(0x1E,0x2E,0x45), C_BOX]
for ri, row in enumerate(rows):
    ty = 2.1 + ri * row_h
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(sl, cx, ty, cw-0.05, row_h-0.03, row_colors[ri % 2])
        col = C_GREEN if "✓" in cell else (RGBColor(0xFF,0x80,0x80) if "✗" in cell else C_WHITE)
        add_text(sl, cell, cx+0.05, ty+0.08, cw-0.1, row_h-0.1,
                 size=Pt(13), color=col, align=PP_ALIGN.CENTER)

# Speed
add_rect(sl, 0.5, 3.85, 7.15, 1.0, C_BOX)
add_text(sl, "辨識速度（攝影機 2）：0.18 – 0.38 秒 / 次",
         0.65, 3.95, 6.8, 0.8, size=Pt(14), color=C_LIGHT)

# Right side: current status (honest)
add_rect(sl, 7.8, 1.0, 5.3, 2.3, C_BOX)
add_rect(sl, 7.8, 1.0, 5.3, 0.45, RGBColor(0x10,0x40,0x60))
add_text(sl, "已改善項目", 7.95, 1.05, 5.0, 0.4,
         size=Pt(13), bold=True, color=C_GREEN)
add_text(sl,
    "✓ 陌生人偵測：攝影機 2 獨立偵測邏輯\n"
    "✓ Grace period 保護剛進入的成員\n"
    "✓ 訓練照片改用攝影機現場拍攝",
    7.95, 1.55, 5.0, 1.6, size=Pt(12), color=C_LIGHT)

add_rect(sl, 7.8, 3.55, 5.3, 3.65, C_BOX)
add_rect(sl, 7.8, 3.55, 5.3, 0.45, RGBColor(0x3A,0x15,0x15))
add_text(sl, "目前仍存在的問題", 7.95, 3.6, 5.0, 0.4,
         size=Pt(13), bold=True, color=RGBColor(0xFF,0x80,0x80))
add_text(sl,
    "⚠ 辨識仍不穩定，低光或角度差時\n  信心值下降導致誤判\n"
    "⚠ 網頁在多攝影機同時串流時\n  偶有延遲或當機\n"
    "⚠ 盲區問題尚未完全解決\n  偶發誤 EXIT",
    7.95, 4.1, 5.0, 2.9, size=Pt(12), color=RGBColor(0xFF,0xC0,0xC0))

# ══════════════════════════════════════════════════════════════════════════════
# Slide 8 — 成果（水耕養殖）
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT)
add_text(sl, "05  研究成果 — 水耕養殖監控", 0.5, 0.2, 12, 0.6,
         size=Pt(28), bold=True, color=C_WHITE)
add_line(sl, 0.5, 0.85, 12.33)

sensors = [
    ("pH", "6.5 – 7.5"),
    ("溫度 (°C)", "22 – 27"),
    ("溶氧量 (mg/L)", "6.0 – 8.0"),
    ("氨氮 (mg/L)", "< 0.5"),
    ("濁度 (NTU)", "< 20"),
    ("電導率 (μS/cm)", "依需求設定"),
]
add_text(sl, "監控指標（3 座養殖池 × 6 項水質參數）",
         0.5, 1.0, 12, 0.45, size=Pt(15), bold=True, color=C_ACCENT)

for i, (name, rng) in enumerate(sensors):
    lx = 0.5 + (i % 3) * 4.2
    ty = 1.55 + (i // 3) * 1.2
    add_rect(sl, lx, ty, 4.0, 1.1, C_BOX)
    add_text(sl, name, lx+0.15, ty+0.07, 3.7, 0.45,
             size=Pt(14), bold=True, color=C_WHITE)
    add_text(sl, f"正常範圍：{rng}", lx+0.15, ty+0.57, 3.7, 0.4,
             size=Pt(12), color=C_LIGHT)

alerts = [
    ("👁  Visual", "數值接近門檻 15%\n卡片顏色警示"),
    ("🟡  Alert 1", "超標持續 ≥ 30 分鐘\nTelegram 通知"),
    ("🔴  Alert 2", "超標持續 ≥ 2 小時\nTelegram 升級警報"),
    ("🔌  Fault",   "回傳值持續為 0\n感測器故障通知"),
]
add_text(sl, "四級警報機制", 0.5, 4.15, 12, 0.45,
         size=Pt(15), bold=True, color=C_ACCENT)
for i, (lvl, desc) in enumerate(alerts):
    lx = 0.5 + i * 3.2
    add_rect(sl, lx, 4.65, 3.0, 2.5, C_BOX)
    add_text(sl, lvl,  lx+0.15, 4.72, 2.7, 0.5,
             size=Pt(14), bold=True, color=C_YELLOW)
    add_text(sl, desc, lx+0.15, 5.3,  2.7, 1.7,
             size=Pt(13), color=C_LIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 9 — 結論
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT)
add_text(sl, "06  結論與未來展望", 0.5, 0.2, 12, 0.6,
         size=Pt(28), bold=True, color=C_WHITE)
add_line(sl, 0.5, 0.85, 12.33)

conclusions = [
    ("技術可行性驗證",
     "低成本 USB 攝影機 + 開源 AI（ArcFace）\n可達實用等級辨識（信心值 0.63–0.79）"),
    ("資料品質勝於數量",
     "與部署環境一致的訓練資料\n比大量低品質資料更有效"),
    ("系統整合效益",
     "人員管制 + 水耕監控整合單一儀表板\n降低管理複雜度與人力成本"),
]
for i, (title, body) in enumerate(conclusions):
    lx = 0.5 + i * 4.2
    add_rect(sl, lx, 1.05, 4.0, 3.2, C_BOX)
    add_rect(sl, lx, 1.05, 4.0, 0.5, C_ACCENT)
    add_text(sl, title, lx+0.15, 1.08, 3.7, 0.45,
             size=Pt(14), bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text(sl, body,  lx+0.15, 1.65, 3.7, 2.4,
             size=Pt(13), color=C_WHITE)

limits = [
    "USB 頻寬限制 → 改用 PoE 網路攝影機",
    "攝影機盲區 → 增加鏡頭數量或調整角度",
    "低光源辨識率下降 → 補光或換低光鏡頭",
    "擴充至 Edge 設備（Jetson / Pi）",
]
future = [
    "擴大人員資料庫測試穩定性",
    "加入出勤統計報表功能",
    "結合 AI 趨勢預測水質異常",
    "開發手機 App 版儀表板",
]

add_bullet_box(sl, limits,  0.5, 4.5, 6.1, 2.75,
               title="已知限制與改善方向", icon="→", item_color=C_LIGHT)
add_bullet_box(sl, future, 6.8, 4.5, 6.1, 2.75,
               title="未來功能擴充",       icon="★", item_color=C_GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# Slide 10 — Thank you / Q&A
# ══════════════════════════════════════════════════════════════════════════════
sl = blank_slide(prs)
fill_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.12, C_ACCENT)
add_rect(sl, 0, 7.38, 13.33, 0.12, C_ACCENT)
add_rect(sl, 0, 0, 0.35, 7.5, C_BOX)
add_rect(sl, 0, 2.5, 0.35, 2.5, C_ACCENT)

add_text(sl, "感謝聆聽", 0.7, 2.0, 12, 1.2,
         size=Pt(52), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "Q & A", 0.7, 3.25, 12, 0.8,
         size=Pt(36), bold=False, color=C_ACCENT, align=PP_ALIGN.CENTER)
add_line(sl, 2.5, 4.2, 8.33)
add_text(sl,
    "智慧實驗室人員管制與水耕養殖監控系統\n"
    "YOLOv8  ·  InsightFace ArcFace  ·  Python Flask  ·  IoT Sensors",
    0.7, 4.5, 12, 1.0,
    size=Pt(14), color=C_GRAY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = r"c:\Users\ITM_Student_06\Desktop\PROJECT-2026-main\專題報告競賽.pptx"
prs.save(out)
print(f"Saved: {out}")
